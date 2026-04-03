import streamlit as st
from google import genai
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from fpdf import FPDF
import io
import os
import re
import time
from PIL import Image
import pytesseract

# --- 1. 페이지 설정 및 데이터 초기화 ---
_ = st.set_page_config(page_title="GapFinder v2 (v24.7)", layout="wide")

states = ['brand_analysis', 'brand_insight', 'comp_analysis', 'consumer_data', 'consumer_analysis', 'final_report']
for key in states:
    if key not in st.session_state:
        st.session_state[key] = "" 

if 'consumer_data' not in st.session_state or not isinstance(st.session_state['consumer_data'], list):
    st.session_state['consumer_data'] = []

# --- 2. 사이드바 설정 ---
with st.sidebar:
    st.header("🔑 GapFinder v2 설정")
    gemini_key = st.text_input("1. Gemini API Key", type="password")
    serper_key = st.text_input("2. Serper API Key", type="password")
    
    st.divider()
    st.subheader("🇳 네이버 검색 API (선택)")
    nav_id = st.text_input("Naver Client ID")
    nav_pw = st.text_input("Naver Client Secret", type="password")
    
    _ = st.divider()
    menu = st.radio("분석 프로세스 단계", [
        "1단계. 브랜드 분석 (Thesis)", 
        "2단계. 경쟁사 분석 (Competitor)", 
        "3단계. 소비자 분석 (Evidence)", 
        "4단계. 통합 전략 리포트 (Synthesis)"
    ])

# --- 3. 유틸리티 함수 ---

def clean_for_pdf(text):
    if not text: return ""
    text = str(text).replace('\u200b', '').replace('\ufeff', '').replace('|', ' ')
    return re.sub(r'[^\u0000-\u007f\uac00-\ud7af\x20-\x7E\s\n\.,!\?\(\)\[\]:;\"\'\-]', '', text)

def extract_all_content(files=None, url=""):
    text = ""
    if files:
        for f in files:
            try:
                if f.name.endswith(".pdf"): text += "\n".join([p.extract_text() for p in PdfReader(f).pages if p.extract_text()])
                elif f.name.endswith(".pptx"): text += "\n".join([s.text for slide in Presentation(f).slides for s in slide.shapes if hasattr(s, "text")])
                elif f.name.lower().endswith((".png", ".jpg", ".jpeg")):
                    text += f"\n[이미지 {f.name} OCR 데이터]:\n" + pytesseract.image_to_string(Image.open(f), lang='kor+eng')
            except: pass
    if url:
        try:
            res = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'}, timeout=5)
            soup = BeautifulSoup(res.text, 'html.parser')
            for s in soup(['script', 'style']): s.decompose()
            text += f"\n[참조 URL: {url}]\n{soup.get_text()[:4000]}"
        except: pass
    return text

def run_ai(data, step, insight="", brand_ctx="", consumer_raw=""):
    """지능형 백오프 재시도 로직이 적용된 전략 엔진 """
    if not gemini_key: return "⚠️ API Key가 필요합니다."
    
    client = genai.Client(api_key=gemini_key)
    p_base = """인사말은 생략하고 15년 차 수석 전략 기획자의 전문적인 비즈니스 언어를 사용하세요. 
    모든 분석은 [기능 / 가격 / 서비스 / 디자인] 4대 카테고리를 준수하세요.\n\n"""
    
    prompts = {
        "brand": f"{p_base}[Thesis] 자사 브랜드 분석: VOI(가치) 및 PDI(리스크)를 10점 만점 지표로 산출하세요. 인사이트: {insight}",
        "comp": f"{p_base}[Competitor] 경쟁 우위 분석: 경쟁사의 전략적 사각지대를 자사({brand_ctx[:300]}) 관점에서 분석하세요.",
        "consumer": f"{p_base}[Evidence] 소비자 여론 분석: 대량 보이스를 카테고리별로 정량 분류하여 Gap을 도출하세요.",
        "final": f"{p_base}[Synthesis] Victory Strategy: Language Gap Matrix, Strategic Foundations, Target-Action Table을 포함한 한 문장 필승 전략 정의. 인사이트: {insight}\n데이터: {consumer_raw[:8000]}"
    }

    # [핵심 업데이트] 지능형 재시도 루프
    max_retries = 3
    for attempt in range(max_retries):
        try:
            res = client.models.generate_content(model="gemini-3-flash-preview", contents=prompts[step] + "\n\n데이터:\n" + str(data)[:25000])
            return res.text
        except Exception as e:
            err_msg = str(e).lower()
            if ("503" in err_msg or "high demand" in err_msg) and attempt < max_retries - 1:
                wait_time = (attempt + 1) * 7 # 7초, 14초 순차적으로 대기 시간 증가 
                st.warning(f"🔄 구글 서버가 혼잡합니다. {wait_time}초 후 재시도합니다... (시도 {attempt + 1}/{max_retries})")
                time.sleep(wait_time)
                continue
            return f"⚠️ 서버 과부하로 분석이 중단되었습니다. 잠시 후(약 1~2분 뒤) '다시 시도' 버튼을 눌러주세요. ({e})"

# --- 4. PDF 및 UI 로직 ---
class SafePDF(FPDF):
    def __init__(self):
        super().__init__()
        f_reg, f_bold = "NanumGothic.ttf", "NanumGothicBold.ttf"
        if os.path.exists(f_reg):
            self.add_font('NG', '', f_reg); self.add_font('NG', 'B', f_bold); self.fn = 'NG'
        else: self.fn = 'Arial'
        _ = self.set_auto_page_break(auto=True, margin=20)
        _ = self.set_margins(20, 20, 20)
    def write_section(self, title, content):
        if not content: return
        self.add_page(); self.set_font(self.fn, 'B', 16); self.set_text_color(0, 51, 102)
        self.cell(170, 15, txt=title, ln=True, align='C'); self.ln(5)
        self.set_font(self.fn, '', 10.5); self.set_text_color(50, 50, 50)
        self.multi_cell(170, 7, txt=clean_for_pdf(content))

# --- 5. 단계별 실행 ---

if menu == "1단계. 브랜드 분석 (Thesis)":
    st.title("🏢 1단계. 브랜드 가치 및 리스크 분석")
    b_f = st.file_uploader("자사 자료 (PDF/PPTX/이미지)", accept_multiple_files=True)
    b_u = st.text_input("자사 웹사이트 URL")
    st.session_state['brand_insight'] = st.text_area("💡 실무/운영 인사이트", value=st.session_state['brand_insight'] if st.session_state['brand_insight'] != "" else "")
    if st.button("브랜드 분석 시작"):
        with st.spinner("서버 상태 확인 및 데이터 추출 중..."):
            st.session_state['brand_analysis'] = run_ai(extract_all_content(b_f, b_u), "brand", st.session_state['brand_insight'])
            _ = st.rerun()
    if st.session_state['brand_analysis']: st.markdown(st.session_state['brand_analysis'])

elif menu == "2단계. 경쟁사 분석 (Competitor)":
    st.title("⚔️ 2단계. 경쟁 구도 분석 (최대 3개)")
    c_f = st.file_uploader("경쟁사 자료", accept_multiple_files=True)
    col1, col2, col3 = st.columns(3)
    with col1: c1n = st.text_input("경쟁사 1"); c1u = st.text_input("URL 1")
    with col2: c2n = st.text_input("경쟁사 2"); c2u = st.text_input("URL 2")
    with col3: c3n = st.text_input("경쟁사 3"); c3u = st.text_input("URL 3")
    if st.button("경쟁 분석 시작"):
        with st.spinner("경쟁 브랜드 틈새 분석 중..."):
            all_c = extract_all_content(c_f)
            for n, u in [(c1n, c1u), (c2n, c2u), (c3n, c3u)]:
                if n:
                    res = requests.post("https://google.serper.dev/search", headers={'X-API-KEY': serper_key}, json={"q": f"{n} 실사용 불만 후기", "num": 30, "gl": "kr", "hl": "ko"}).json()
                    all_c += f"\n[{n} 데이터]\n" + "\n".join([r.get('snippet', '') for r in res.get('organic', [])])
                    if u: all_c += extract_all_content(url=u)
            # 자사 분석 컨텍스트를 넘겨서 비교 분석 수행
            st.session_state['comp_analysis'] = run_ai(all_c, "comp", brand_ctx=st.session_state['brand_analysis'])
            _ = st.rerun()
    if st.session_state['comp_analysis']: st.markdown(st.session_state['comp_analysis'])

elif menu == "3단계. 소비자 분석 (Evidence)":
    st.title("👥 3단계. 소비자 보이스 정량 분류")
    kw = st.text_input("분석 키워드 입력", value="")
    if st.button("데이터 수집 및 분석 시작"):
        if not kw: st.warning("키워드를 입력해주세요.")
        else:
            with st.spinner("대량 데이터 수집 중..."):
                all_r = []
                for k in [x.strip() for x in kw.split(",")]:
                    res = requests.post("https://google.serper.dev/search", headers={'X-API-KEY': serper_key}, json={"q": f"{k}", "num": 50, "gl": "kr", "hl": "ko"}).json()
                    for r in res.get('organic', []):
                        tag = "[🔴유튜브]" if "youtube" in r.get('link', '') else "[🔵네이버]"
                        all_r.append(f"{tag} {r.get('title')}: {r.get('snippet')}")
                if all_r:
                    st.session_state['consumer_data'] = all_r
                    st.session_state['consumer_analysis'] = run_ai("\n".join(all_r), "consumer")
                    _ = st.rerun()
    if st.session_state['consumer_analysis']: 
        st.markdown(st.session_state['consumer_analysis'])
        st.divider()
        col_c1, col_c2 = st.columns([2, 1])
        with col_c1:
            with st.expander(f"📝 원본 데이터 확인"):
                for line in st.session_state['consumer_data']: st.write(line)
        with col_c2:
            st.download_button(f"📥 에비던스 팩 다운로드", data="\n".join(st.session_state['consumer_data']), file_name="Evidence_Full.txt")

elif menu == "4단계. 통합 전략 리포트 (Synthesis)":
    st.title("🧠 4단계. 최종 Victory Strategy")
    if st.button("🚀 최종 리포트 생성"):
        with st.spinner("전략 합성 중..."):
            comb = f"자사:{st.session_state['brand_analysis']}\n경쟁:{st.session_state['comp_analysis']}\n소비자:{st.session_state['consumer_analysis']}"
            st.session_state['final_report'] = run_ai(comb, "final", st.session_state['brand_insight'], consumer_raw=str(st.session_state['consumer_data']))
            _ = st.rerun()
    if st.session_state['final_report']:
        st.markdown(st.session_state['final_report'])
        st.divider()
        pdf = SafePDF()
        pdf.write_section("1. STRATEGIC THESIS", st.session_state['brand_analysis'])
        pdf.write_section("2. COMPETITIVE GAP", st.session_state['comp_analysis'])
        pdf.write_section("3. CONSUMER INSIGHT", st.session_state['consumer_analysis'])
        pdf.write_section("4. VICTORY STRATEGY v24.7", st.session_state['final_report'])
        st.download_button("📥 통합 리포트 PDF 다운로드", data=bytes(pdf.output()), file_name="GapFinder_v24_Final.pdf", mime="application/pdf")
